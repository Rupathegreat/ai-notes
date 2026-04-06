from fastapi import FastAPI, APIRouter, HTTPException, UploadFile, File, Request, Header, Depends
from fastapi.responses import JSONResponse
from dotenv import load_dotenv
from starlette.middleware.cors import CORSMiddleware
from motor.motor_asyncio import AsyncIOMotorClient
from pydantic import BaseModel, Field, ConfigDict
from typing import List, Optional, Dict, Any
from datetime import datetime, timezone, timedelta
import os
import logging
from pathlib import Path
import uuid
import aiofiles
import asyncio
import subprocess
from pypdf import PdfReader
from pptx import Presentation
import io
import httpx
from emergentintegrations.llm.chat import LlmChat, UserMessage
from emergentintegrations.llm.openai import OpenAISpeechToText

# Load environment variables
ROOT_DIR = Path(__file__).parent
load_dotenv(ROOT_DIR / '.env')

# MongoDB connection
mongo_url = os.environ['MONGO_URL']
client = AsyncIOMotorClient(mongo_url)
db = client[os.environ['DB_NAME']]

# Create the main app
app = FastAPI()
api_router = APIRouter(prefix="/api")

# Upload directory
UPLOAD_DIR = Path(os.environ.get('UPLOAD_DIR', '/app/backend/uploads'))
UPLOAD_DIR.mkdir(exist_ok=True, parents=True)

# Configure logging
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(name)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

# ==================== MODELS ====================

class User(BaseModel):
    model_config = ConfigDict(extra="ignore")
    user_id: str
    email: str
    name: str
    picture: Optional[str] = None
    language_preference: str = "en"
    theme_preference: str = "light"
    created_at: datetime

class UserSession(BaseModel):
    model_config = ConfigDict(extra="ignore")
    session_token: str
    user_id: str
    expires_at: datetime
    created_at: datetime

class Lecture(BaseModel):
    model_config = ConfigDict(extra="ignore")
    lecture_id: str
    user_id: str
    title: str
    file_name: str
    file_path: str
    file_type: str
    status: str  # 'processing', 'completed', 'failed'
    error_message: Optional[str] = None
    created_at: datetime

class Note(BaseModel):
    model_config = ConfigDict(extra="ignore")
    note_id: str
    lecture_id: str
    user_id: str
    title: str
    summary: str
    key_concepts: List[str]
    important_points: List[str]
    definitions: List[Dict[str, str]]
    keywords: List[str]
    faq: List[Dict[str, str]]
    quiz: List[Dict[str, Any]]
    flowchart: str
    language: str = "en"
    created_at: datetime

class ChatMessage(BaseModel):
    model_config = ConfigDict(extra="ignore")
    message_id: str
    lecture_id: str
    user_id: str
    role: str  # 'user' or 'assistant'
    content: str
    timestamp: datetime

class UserPreferences(BaseModel):
    language: Optional[str] = "en"
    theme: Optional[str] = "light"

# Request/Response models
class SessionCallbackRequest(BaseModel):
    session_id: str

class ChatRequest(BaseModel):
    message: str

# ==================== AUTHENTICATION ====================

async def get_current_user(
    request: Request,
    authorization: Optional[str] = Header(None)
) -> User:
    """
    Authenticator helper - checks session_token from cookies first, then Authorization header.
    WARNING: Don't use FastAPI's HTTPAuthorizationCredentials dependency - it breaks cookie auth.
    """
    session_token = None
    
    # Check cookies first
    session_token = request.cookies.get("session_token")
    
    # Fallback to Authorization header
    if not session_token and authorization:
        if authorization.startswith("Bearer "):
            session_token = authorization.replace("Bearer ", "")
        else:
            session_token = authorization
    
    if not session_token:
        raise HTTPException(status_code=401, detail="Not authenticated")
    
    # Find session in database
    session_doc = await db.user_sessions.find_one(
        {"session_token": session_token},
        {"_id": 0}
    )
    
    if not session_doc:
        raise HTTPException(status_code=401, detail="Invalid session")
    
    # Check expiry
    expires_at = session_doc["expires_at"]
    if isinstance(expires_at, str):
        expires_at = datetime.fromisoformat(expires_at)
    if expires_at.tzinfo is None:
        expires_at = expires_at.replace(tzinfo=timezone.utc)
    if expires_at < datetime.now(timezone.utc):
        raise HTTPException(status_code=401, detail="Session expired")
    
    # Get user
    user_doc = await db.users.find_one(
        {"user_id": session_doc["user_id"]},
        {"_id": 0}
    )
    
    if not user_doc:
        raise HTTPException(status_code=404, detail="User not found")
    
    return User(**user_doc)

@api_router.post("/auth/google-callback")
async def google_callback(request: SessionCallbackRequest):
    """Exchange session_id for session_token via Emergent Auth"""
    try:
        async with httpx.AsyncClient() as client:
            response = await client.get(
                "https://demobackend.emergentagent.com/auth/v1/env/oauth/session-data",
                headers={"X-Session-ID": request.session_id},
                timeout=10.0
            )
            response.raise_for_status()
            data = response.json()
        
        # Create or update user
        user_id = f"user_{uuid.uuid4().hex[:12]}"
        existing_user = await db.users.find_one({"email": data["email"]}, {"_id": 0})
        
        if existing_user:
            user_id = existing_user["user_id"]
            # Update user info
            await db.users.update_one(
                {"user_id": user_id},
                {"$set": {
                    "name": data["name"],
                    "picture": data.get("picture")
                }}
            )
        else:
            # Create new user
            user_doc = {
                "user_id": user_id,
                "email": data["email"],
                "name": data["name"],
                "picture": data.get("picture"),
                "language_preference": "en",
                "theme_preference": "light",
                "created_at": datetime.now(timezone.utc).isoformat()
            }
            await db.users.insert_one(user_doc)
        
        # Store session
        session_doc = {
            "session_token": data["session_token"],
            "user_id": user_id,
            "expires_at": (datetime.now(timezone.utc) + timedelta(days=7)).isoformat(),
            "created_at": datetime.now(timezone.utc).isoformat()
        }
        await db.user_sessions.insert_one(session_doc)
        
        # Get user
        user_doc = await db.users.find_one({"user_id": user_id}, {"_id": 0})
        
        return {
            "session_token": data["session_token"],
            "user": user_doc
        }
    
    except httpx.HTTPError as e:
        logger.error(f"Emergent Auth error: {e}")
        raise HTTPException(status_code=400, detail="Failed to authenticate with Google")
    except Exception as e:
        logger.error(f"Auth error: {e}")
        raise HTTPException(status_code=500, detail="Authentication failed")

@api_router.get("/auth/me")
async def get_me(current_user: User = Depends(get_current_user)):
    """Get current authenticated user"""
    return current_user

@api_router.post("/auth/logout")
async def logout(request: Request, current_user: User = Depends(get_current_user)):
    """Logout user"""
    session_token = request.cookies.get("session_token")
    if session_token:
        await db.user_sessions.delete_one({"session_token": session_token})
    return {"message": "Logged out successfully"}

# ==================== USER PREFERENCES ====================

@api_router.get("/user/preferences")
async def get_preferences(current_user: User = Depends(get_current_user)):
    """Get user preferences"""
    return {
        "language": current_user.language_preference,
        "theme": current_user.theme_preference
    }

@api_router.put("/user/preferences")
async def update_preferences(
    preferences: UserPreferences,
    current_user: User = Depends(get_current_user)
):
    """Update user preferences"""
    update_data = {}
    if preferences.language:
        update_data["language_preference"] = preferences.language
    if preferences.theme:
        update_data["theme_preference"] = preferences.theme
    
    if update_data:
        await db.users.update_one(
            {"user_id": current_user.user_id},
            {"$set": update_data}
        )
    
    return {"message": "Preferences updated successfully"}

# ==================== FILE PROCESSING UTILITIES ====================

async def extract_text_from_pdf(file_path: str) -> str:
    """Extract text from PDF file"""
    try:
        reader = PdfReader(file_path)
        text = ""
        for page in reader.pages:
            text += page.extract_text() + "\n"
        return text.strip()
    except Exception as e:
        logger.error(f"PDF extraction error: {e}")
        raise Exception(f"Failed to extract text from PDF: {str(e)}")

async def extract_text_from_pptx(file_path: str) -> str:
    """Extract text from PPTX file"""
    try:
        prs = Presentation(file_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text.strip()
    except Exception as e:
        logger.error(f"PPTX extraction error: {e}")
        raise Exception(f"Failed to extract text from PPTX: {str(e)}")

async def extract_text_from_txt(file_path: str) -> str:
    """Extract text from TXT file"""
    try:
        async with aiofiles.open(file_path, 'r', encoding='utf-8') as f:
            text = await f.read()
        return text.strip()
    except Exception as e:
        logger.error(f"TXT extraction error: {e}")
        raise Exception(f"Failed to read text file: {str(e)}")

async def transcribe_audio_file(file_path: str) -> str:
    """Transcribe audio/video file using Whisper"""
    try:
        stt = OpenAISpeechToText(api_key=os.getenv("EMERGENT_LLM_KEY"))
        
        with open(file_path, "rb") as audio_file:
            response = await stt.transcribe(
                file=audio_file,
                model="whisper-1",
                response_format="text"
            )
        
        return response.strip()
    except Exception as e:
        logger.error(f"Transcription error: {e}")
        raise Exception(f"Failed to transcribe audio: {str(e)}")

async def download_video_from_url(url: str, output_path: str) -> str:
    """Download video/audio from URL (YouTube, etc.) using yt-dlp"""
    try:
        # Use full path to yt-dlp
        yt_dlp_path = '/root/.venv/bin/yt-dlp'
        
        # Simple and reliable command without cookies
        command = [
            yt_dlp_path,
            '-f', 'bestaudio/best',
            '-x',  # Extract audio
            '--audio-format', 'mp3',
            '--audio-quality', '5',
            '--no-playlist',
            '--max-filesize', '1000M',
            '--no-check-certificates',
            '--prefer-free-formats',
            '--user-agent', 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36',
            '-o', output_path,
            url
        ]
        
        logger.info(f"Downloading from URL: {url}")
        
        # Run in thread pool
        loop = asyncio.get_event_loop()
        result = await loop.run_in_executor(
            None, 
            lambda: subprocess.run(command, capture_output=True, text=True, timeout=900)
        )
        
        # Check for output file
        final_path = output_path if output_path.endswith('.mp3') else f"{output_path}.mp3"
        
        if os.path.exists(final_path):
            file_size = os.path.getsize(final_path)
            logger.info(f"Successfully downloaded to {final_path} ({file_size} bytes)")
            return final_path
        
        # Log full error for debugging
        error_msg = result.stderr or result.stdout or "Unknown error"
        logger.error(f"Download failed: {error_msg[:500]}")
        
        # Return more user-friendly error
        if "403" in error_msg or "Forbidden" in error_msg:
            raise Exception("Video blocked by platform. Try: 1) Direct file upload, 2) Vimeo links, 3) Direct MP4/MP3 URLs")
        elif "not available" in error_msg.lower():
            raise Exception("Video not available or private")
        else:
            raise Exception("Download failed. Try uploading the file directly instead.")
        
    except subprocess.TimeoutExpired:
        raise Exception("Download timed out. Try a shorter video or upload file directly.")
    except Exception as e:
        logger.error(f"Video download error: {e}")
        raise

async def generate_structured_notes(text: str, title: str) -> Dict[str, Any]:
    """Generate structured notes using GPT-5.2"""
    try:
        chat = LlmChat(
            api_key=os.getenv("EMERGENT_LLM_KEY"),
            session_id=f"notes_{uuid.uuid4().hex[:8]}",
            system_message="You are an expert educational content analyzer. Generate structured learning materials from lecture content."
        ).with_model("openai", "gpt-5.2")
        
        # Limit text to 8000 chars for faster processing
        text_sample = text[:8000] if len(text) > 8000 else text
        
        prompt = f"""Analyze this lecture content and generate comprehensive structured notes in JSON format.

Lecture Title: {title}

Content:
{text_sample}

Generate a JSON response with the following structure:
{{
  "title": "string (concise title)",
  "summary": "string (2-3 paragraph summary)",
  "key_concepts": ["concept1", "concept2", ...] (5-10 main concepts),
  "important_points": ["point1", "point2", ...] (8-15 bullet points),
  "definitions": [{{"term": "string", "definition": "string"}}, ...] (5-10 key terms),
  "keywords": ["keyword1", "keyword2", ...] (10-15 keywords),
  "faq": [{{"question": "string", "answer": "string"}}, ...] (5-8 FAQs),
  "quiz": [{{"question": "string", "options": ["A", "B", "C", "D"], "correctAnswerIndex": 0, "explanation": "why this is correct"}}, ...] (5 questions),
  "flowchart": "string (Mermaid flowchart syntax - MUST be a proper flowchart showing process flow)"
}}

IMPORTANT FOR FLOWCHART:
- Use proper Mermaid flowchart syntax
- Start with: flowchart TD
- Use shapes: [] for process, {{}} for decision, [()] for start/end
- Connect with arrows: -->
- Example format:
  flowchart TD
      A[Start] --> B{{Decision?}}
      B -->|Yes| C[Process 1]
      B -->|No| D[Process 2]
      C --> E[End]
      D --> E

Create a flowchart that shows the main flow/process of the lecture topic.

IMPORTANT: Return ONLY valid JSON, no markdown, no explanation."""

        user_message = UserMessage(text=prompt)
        response = await chat.send_message(user_message)
        
        # Parse JSON response
        import json
        
        # Clean response
        response_text = response.strip()
        if response_text.startswith("```json"):
            response_text = response_text[7:]
        if response_text.startswith("```"):
            response_text = response_text[3:]
        if response_text.endswith("```"):
            response_text = response_text[:-3]
        response_text = response_text.strip()
        
        notes_data = json.loads(response_text)
        
        # Validate and fix flowchart if needed
        if not notes_data.get('flowchart') or len(notes_data.get('flowchart', '')) < 20:
            notes_data['flowchart'] = """flowchart TD
    Start([Start Learning]) --> Intro[Introduction to Topic]
    Intro --> Core[Core Concepts]
    Core --> Practice{{Practice Exercises?}}
    Practice -->|Yes| Apply[Apply Knowledge]
    Practice -->|No| Review[Review Material]
    Apply --> Test[Take Quiz]
    Review --> Test
    Test --> End([Complete])"""
        
        return notes_data
    
    except Exception as e:
        logger.error(f"Notes generation error: {e}")
        # Return fallback structure
        return {
            "title": title,
            "summary": "Failed to generate summary. Please try again.",
            "key_concepts": ["Content processing in progress"],
            "important_points": ["Notes generation encountered an error"],
            "definitions": [{"term": "Error", "definition": "Failed to extract definitions"}],
            "keywords": ["error"],
            "faq": [{"question": "Why did this fail?", "answer": "The AI service encountered an error. Please try again."}],
            "quiz": [{"question": "Sample question", "options": ["A", "B", "C", "D"], "correctAnswerIndex": 0, "explanation": "This is a sample explanation"}],
            "flowchart": """flowchart TD
    Start([Start]) --> Process[Processing]
    Process --> Error{{Error Occurred}}
    Error --> Retry[Please Retry]
    Retry --> End([End])"""
        }

async def process_lecture_background(lecture_id: str, file_path: str, file_type: str, title: str, user_id: str, is_url: bool = False):
    """Background task to process lecture and generate notes"""
    try:
        # Update status to processing
        await db.lectures.update_one(
            {"lecture_id": lecture_id},
            {"$set": {"status": "processing"}}
        )
        
        # Extract text based on file type
        text = ""
        
        if is_url:
            # For URLs, download first then transcribe
            text = await transcribe_audio_file(file_path)
        elif file_type == "application/pdf":
            text = await extract_text_from_pdf(file_path)
        elif file_type in ["application/vnd.ms-powerpoint", "application/vnd.openxmlformats-officedocument.presentationml.presentation"]:
            text = await extract_text_from_pptx(file_path)
        elif file_type == "text/plain":
            text = await extract_text_from_txt(file_path)
        elif file_type.startswith("audio/") or file_type.startswith("video/"):
            text = await transcribe_audio_file(file_path)
        else:
            raise Exception(f"Unsupported file type: {file_type}")
        
        if not text or len(text) < 50:
            raise Exception("Insufficient text extracted from file")
        
        # Generate structured notes
        notes_data = await generate_structured_notes(text, title)
        
        # Save notes to database
        note_doc = {
            "note_id": f"note_{uuid.uuid4().hex[:12]}",
            "lecture_id": lecture_id,
            "user_id": user_id,
            "title": notes_data.get("title", title),
            "summary": notes_data.get("summary", ""),
            "key_concepts": notes_data.get("key_concepts", []),
            "important_points": notes_data.get("important_points", []),
            "definitions": notes_data.get("definitions", []),
            "keywords": notes_data.get("keywords", []),
            "faq": notes_data.get("faq", []),
            "quiz": notes_data.get("quiz", []),
            "flowchart": notes_data.get("flowchart", ""),
            "language": "en",
            "created_at": datetime.now(timezone.utc).isoformat()
        }
        await db.notes.insert_one(note_doc)
        
        # Update lecture status
        await db.lectures.update_one(
            {"lecture_id": lecture_id},
            {"$set": {"status": "completed"}}
        )
        
        logger.info(f"Lecture {lecture_id} processed successfully")
        
    except Exception as e:
        logger.error(f"Lecture processing error: {e}")
        await db.lectures.update_one(
            {"lecture_id": lecture_id},
            {"$set": {
                "status": "failed",
                "error_message": str(e)
            }}
        )

# ==================== LECTURE ENDPOINTS ====================

@api_router.post("/lectures/upload")
async def upload_lecture(
    file: UploadFile = File(...),
    title: Optional[str] = None,
    current_user: User = Depends(get_current_user)
):
    """Upload a lecture file"""
    try:
        # Validate file type
        allowed_types = [
            "application/pdf",
            "application/vnd.ms-powerpoint",
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            "text/plain",
            "audio/mpeg", "audio/mp3", "audio/wav", "audio/m4a",
            "video/mp4", "video/mpeg", "video/quicktime"
        ]
        
        if file.content_type not in allowed_types:
            raise HTTPException(status_code=400, detail=f"Unsupported file type: {file.content_type}")
        
        # Generate lecture ID and file path
        lecture_id = f"lecture_{uuid.uuid4().hex[:12]}"
        file_extension = Path(file.filename).suffix
        file_path = UPLOAD_DIR / f"{lecture_id}{file_extension}"
        
        # Save file
        async with aiofiles.open(file_path, 'wb') as f:
            content = await file.read()
            await f.write(content)
        
        # Create lecture document
        lecture_doc = {
            "lecture_id": lecture_id,
            "user_id": current_user.user_id,
            "title": title or file.filename,
            "file_name": file.filename,
            "file_path": str(file_path),
            "file_type": file.content_type,
            "status": "processing",
            "error_message": None,
            "created_at": datetime.now(timezone.utc).isoformat()
        }
        await db.lectures.insert_one(lecture_doc)
        
        # Start background processing
        asyncio.create_task(
            process_lecture_background(
                lecture_id=lecture_id,
                file_path=str(file_path),
                file_type=file.content_type,
                title=title or file.filename,
                user_id=current_user.user_id
            )
        )
        
        return {
            "lecture_id": lecture_id,
            "status": "processing",
            "message": "Lecture uploaded successfully. Processing started."
        }
    
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Upload error: {e}")
        raise HTTPException(status_code=500, detail=f"Upload failed: {str(e)}")

@api_router.get("/lectures/{lecture_id}/status")
async def get_lecture_status(
    lecture_id: str,
    current_user: User = Depends(get_current_user)
):
    """Get lecture processing status"""
    lecture_doc = await db.lectures.find_one(
        {"lecture_id": lecture_id, "user_id": current_user.user_id},
        {"_id": 0}
    )
    
    if not lecture_doc:
        raise HTTPException(status_code=404, detail="Lecture not found")
    
    return {
        "lecture_id": lecture_id,
        "status": lecture_doc["status"],
        "error_message": lecture_doc.get("error_message")
    }

@api_router.get("/lectures")
async def get_lectures(current_user: User = Depends(get_current_user)):
    """Get all lectures for current user"""
    lectures = await db.lectures.find(
        {"user_id": current_user.user_id},
        {"_id": 0}
    ).sort("created_at", -1).to_list(100)
    
    return lectures

@api_router.get("/lectures/analytics/stats")
async def get_analytics(current_user: User = Depends(get_current_user)):
    """Get analytics for current user's lectures"""
    try:
        # Get all lectures
        all_lectures = await db.lectures.find(
            {"user_id": current_user.user_id},
            {"_id": 0}
        ).to_list(1000)
        
        # Calculate stats
        total_lectures = len(all_lectures)
        completed_lectures = sum(1 for lec in all_lectures if lec.get("status") == "completed")
        processing_lectures = sum(1 for lec in all_lectures if lec.get("status") == "processing")
        failed_lectures = sum(1 for lec in all_lectures if lec.get("status") == "failed")
        
        # Recent activity (last 5)
        recent_lectures = sorted(all_lectures, key=lambda x: x.get("created_at", ""), reverse=True)[:5]
        
        return {
            "total_lectures": total_lectures,
            "completed": completed_lectures,
            "processing": processing_lectures,
            "failed": failed_lectures,
            "recent_activity": recent_lectures
        }
    except Exception as e:
        logger.error(f"Analytics error: {e}")
        return {
            "total_lectures": 0,
            "completed": 0,
            "processing": 0,
            "failed": 0,
            "recent_activity": []
        }

@api_router.get("/lectures/{lecture_id}")
async def get_lecture(
    lecture_id: str,
    current_user: User = Depends(get_current_user)
):
    """Get lecture details"""
    lecture_doc = await db.lectures.find_one(
        {"lecture_id": lecture_id, "user_id": current_user.user_id},
        {"_id": 0}
    )
    
    if not lecture_doc:
        raise HTTPException(status_code=404, detail="Lecture not found")
    
    return lecture_doc

class UrlUploadRequest(BaseModel):
    url: str

@api_router.post("/lectures/upload-url")
async def upload_lecture_url(
    request: UrlUploadRequest,
    current_user: User = Depends(get_current_user)
):
    """Upload a lecture from URL (YouTube, Vimeo, etc.)"""
    try:
        url = request.url.strip()
        
        if not url:
            raise HTTPException(status_code=400, detail="URL is required")
        
        # Generate lecture ID and file path
        lecture_id = f"lecture_{uuid.uuid4().hex[:12]}"
        audio_path = UPLOAD_DIR / f"{lecture_id}"
        
        # Extract title from URL
        title = url.split('/')[-1][:50] or "Video Lecture"
        
        # Create lecture document
        lecture_doc = {
            "lecture_id": lecture_id,
            "user_id": current_user.user_id,
            "title": title,
            "file_name": url,
            "file_path": str(audio_path),
            "file_type": "video/url",
            "status": "processing",
            "error_message": None,
            "created_at": datetime.now(timezone.utc).isoformat()
        }
        await db.lectures.insert_one(lecture_doc)
        
        # Download video and process in background
        async def download_and_process():
            try:
                # Download video/audio
                downloaded_path = await download_video_from_url(url, str(audio_path))
                
                # Update file path
                await db.lectures.update_one(
                    {"lecture_id": lecture_id},
                    {"$set": {"file_path": downloaded_path}}
                )
                
                # Process lecture
                await process_lecture_background(
                    lecture_id=lecture_id,
                    file_path=downloaded_path,
                    file_type="audio/mp3",
                    title=title,
                    user_id=current_user.user_id,
                    is_url=True
                )
            except Exception as e:
                logger.error(f"URL processing error: {e}")
                await db.lectures.update_one(
                    {"lecture_id": lecture_id},
                    {"$set": {
                        "status": "failed",
                        "error_message": str(e)
                    }}
                )
        
        # Start background task
        asyncio.create_task(download_and_process())
        
        return {
            "lecture_id": lecture_id,
            "status": "processing",
            "message": "URL submitted successfully. Downloading and processing started."
        }
    
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"URL upload error: {e}")
        raise HTTPException(status_code=500, detail=f"URL upload failed: {str(e)}")

# ==================== NOTES ENDPOINTS ====================

@api_router.get("/notes/{lecture_id}")
async def get_notes(
    lecture_id: str,
    current_user: User = Depends(get_current_user)
):
    """Get generated notes for a lecture"""
    # Verify lecture belongs to user
    lecture_doc = await db.lectures.find_one(
        {"lecture_id": lecture_id, "user_id": current_user.user_id},
        {"_id": 0}
    )
    
    if not lecture_doc:
        raise HTTPException(status_code=404, detail="Lecture not found")
    
    # Get notes
    notes_doc = await db.notes.find_one(
        {"lecture_id": lecture_id},
        {"_id": 0}
    )
    
    if not notes_doc:
        raise HTTPException(status_code=404, detail="Notes not generated yet")
    
    return notes_doc

# ==================== CHATBOT ENDPOINTS ====================

@api_router.post("/chat/{lecture_id}")
async def send_chat_message(
    lecture_id: str,
    chat_request: ChatRequest,
    current_user: User = Depends(get_current_user)
):
    """Send message to chatbot"""
    try:
        # Verify lecture belongs to user
        lecture_doc = await db.lectures.find_one(
            {"lecture_id": lecture_id, "user_id": current_user.user_id},
            {"_id": 0}
        )
        
        if not lecture_doc:
            raise HTTPException(status_code=404, detail="Lecture not found")
        
        # Get notes for context
        notes_doc = await db.notes.find_one(
            {"lecture_id": lecture_id},
            {"_id": 0}
        )
        
        # Build comprehensive context
        context = ""
        if notes_doc:
            # Include all important information
            concepts = notes_doc.get('key_concepts', [])
            points = notes_doc.get('important_points', [])
            definitions = notes_doc.get('definitions', [])
            faq = notes_doc.get('faq', [])
            
            context = f"""You have access to comprehensive lecture notes. Here's what you know:

LECTURE TITLE: {notes_doc.get('title', 'Unknown')}

SUMMARY:
{notes_doc.get('summary', 'No summary available')}

KEY CONCEPTS:
{chr(10).join(f'- {concept}' for concept in concepts[:10])}

IMPORTANT POINTS:
{chr(10).join(f'- {point}' for point in points[:15])}

DEFINITIONS:
{chr(10).join(f'- {d.get("term", "")}: {d.get("definition", "")}' for d in definitions[:10])}

FREQUENTLY ASKED QUESTIONS:
{chr(10).join(f'Q: {q.get("question", "")}{chr(10)}A: {q.get("answer", "")}' for q in faq[:5])}

KEYWORDS: {', '.join(notes_doc.get('keywords', [])[:20])}
"""
        else:
            context = "No lecture notes available yet. The lecture may still be processing."
        
        # Get chat history for continuity
        chat_history = await db.chat_messages.find(
            {"lecture_id": lecture_id, "user_id": current_user.user_id},
            {"_id": 0}
        ).sort("timestamp", 1).limit(10).to_list(10)
        
        # Build conversation history
        conversation_context = ""
        if chat_history:
            conversation_context = "\n\nPrevious conversation:\n"
            for msg in chat_history[-5:]:  # Last 5 messages
                role = "Student" if msg['role'] == 'user' else "You"
                conversation_context += f"{role}: {msg['content']}\n"
        
        # Initialize chat with better system message
        chat = LlmChat(
            api_key=os.getenv("EMERGENT_LLM_KEY"),
            session_id=f"chat_{lecture_id}_{current_user.user_id}",
            system_message=f"""You are an expert AI tutor and teaching assistant. Your role is to help students understand this lecture material deeply.

{context}
{conversation_context}

GUIDELINES:
- Provide detailed, clear explanations using the lecture content
- Use examples and analogies to clarify concepts
- Break down complex topics into simpler parts
- If asked about something not in the lecture, clearly state that and provide general knowledge if helpful
- Be encouraging and supportive
- Use bullet points and formatting for clarity
- Reference specific concepts, definitions, or points from the lecture notes when relevant
- If a student seems confused, ask clarifying questions"""
        ).with_model("openai", "gpt-5.2")
        
        # Send message
        user_message = UserMessage(text=chat_request.message)
        response = await chat.send_message(user_message)
        
        # Save user message
        user_msg_doc = {
            "message_id": f"msg_{uuid.uuid4().hex[:12]}",
            "lecture_id": lecture_id,
            "user_id": current_user.user_id,
            "role": "user",
            "content": chat_request.message,
            "timestamp": datetime.now(timezone.utc).isoformat()
        }
        await db.chat_messages.insert_one(user_msg_doc)
        
        # Save assistant response
        assistant_msg_doc = {
            "message_id": f"msg_{uuid.uuid4().hex[:12]}",
            "lecture_id": lecture_id,
            "user_id": current_user.user_id,
            "role": "assistant",
            "content": response,
            "timestamp": datetime.now(timezone.utc).isoformat()
        }
        await db.chat_messages.insert_one(assistant_msg_doc)
        
        return {
            "message": response
        }
    
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Chat error: {e}")
        raise HTTPException(status_code=500, detail=f"Chat failed: {str(e)}")

@api_router.get("/chat/{lecture_id}/history")
async def get_chat_history(
    lecture_id: str,
    current_user: User = Depends(get_current_user)
):
    """Get chat history for a lecture"""
    # Verify lecture belongs to user
    lecture_doc = await db.lectures.find_one(
        {"lecture_id": lecture_id, "user_id": current_user.user_id},
        {"_id": 0}
    )
    
    if not lecture_doc:
        raise HTTPException(status_code=404, detail="Lecture not found")
    
    # Get chat history
    messages = await db.chat_messages.find(
        {"lecture_id": lecture_id, "user_id": current_user.user_id},
        {"_id": 0}
    ).sort("timestamp", 1).to_list(100)
    
    return messages

# ==================== ROOT ENDPOINT ====================

@api_router.get("/")
async def root():
    return {"message": "Smart Notes API", "status": "running"}

# Include router
app.include_router(api_router)

# CORS
app.add_middleware(
    CORSMiddleware,
    allow_credentials=True,
    allow_origins=os.environ.get('CORS_ORIGINS', '*').split(','),
    allow_methods=["*"],
    allow_headers=["*"],
)

@app.on_event("shutdown")
async def shutdown_db_client():
    client.close()
